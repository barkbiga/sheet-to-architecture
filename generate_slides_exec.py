
#!/usr/bin/env python3
"""Génère un deck PPTX exécutif contenant overview.png + KPI.
Requiert python-pptx, pandas.
"""
import argparse, pathlib, pandas as pd, subprocess, shutil
from pptx import Presentation
from pptx.util import Inches

def create_ppt(xlsx, diagrams_dir, out_file):
    ctx = pd.ExcelFile(xlsx)
    project = ctx.parse('Project').iloc[0].to_dict()
    objectives = ctx.parse('Objectives')
    prs = Presentation()
    # titre
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = project.get('Context','Projet')
    slide.placeholders[1].text = 'Sponsor: ' + project.get('Sponsor','')
    # Overview avec image
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = 'Vue d’ensemble'
    img = diagrams_dir / 'overview.png'
    if img.exists():
        slide2.shapes.add_picture(str(img), Inches(1), Inches(2), width=Inches(8))
    else:
        slide2.shapes.placeholders[1].text_frame.text = 'Diagramme overview non généré.'
    # KPI slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = 'KPI clés'
    tf = slide3.shapes.placeholders[1].text_frame
    tf.clear()
    for _, row in objectives.iterrows():
        p = tf.add_paragraph()
        p.text = f"{row['KPI']} ⇒ {row['TargetValue']} {row['Unit']}"
    prs.save(out_file)
    print('✅ Slide deck créé :', out_file)

if __name__ == '__main__':
    import argparse, pathlib
    ap = argparse.ArgumentParser()
    ap.add_argument('-i','--input',default='petstore_archi_v3.xlsx')
    ap.add_argument('-o','--output',default='executive_deck.pptx')
    ap.add_argument('-d','--diagrams',default='diagrams')
    args = ap.parse_args()
    create_ppt(pathlib.Path(args.input), pathlib.Path(args.diagrams), pathlib.Path(args.output))

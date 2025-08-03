
#!/usr/bin/env python3
"""generate_presentation.py
Génère une présentation PowerPoint (.pptx) à partir du classeur d’architecture Excel.

Usage:
    python generate_presentation.py -i petstore_archi_v3.xlsx -o presentation.pptx
Nécessite: pandas, python-pptx (pip install python-pptx)
"""

import argparse, pathlib, pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

TITLE_SLIDE_LAYOUT = 0
TITLE_AND_CONTENT_LAYOUT = 1

def read_ctx(xlsx: pathlib.Path):
    wb = pd.ExcelFile(xlsx)
    ctx = {s: wb.parse(s) for s in wb.sheet_names}
    project = wb.parse('Project').iloc[0].to_dict()
    return ctx, project

def add_title_slide(prs: Presentation, project):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE_LAYOUT])
    slide.shapes.title.text = project.get('Context', 'Projet')
    subtitle = slide.placeholders[1]
    sponsor = project.get('Sponsor', '')
    subtitle.text = f"Sponsor: {sponsor}"

def add_objectives_slide(prs: Presentation, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = "Objectifs & KPI"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for _, row in ctx['Objectives'].iterrows():
        p = body.add_paragraph()
        p.level = 0
        p.text = f"{row['Label']} — {row['KPI']} (cible {row['TargetValue']} {row['Unit']})"

def add_process_slide(prs: Presentation, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = "Processus métier"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for _, row in ctx['BusinessProcesses'].iterrows():
        p = body.add_paragraph()
        p.text = f"[{row['ID']}] {row['Name']}"
        if pd.notna(row.get('Description')):
            p2 = body.add_paragraph()
            p2.level = 1
            p2.text = f"Description: {row['Description']}"
        if pd.notna(row.get('Actors')):
            p3 = body.add_paragraph()
            p3.level = 1
            p3.text = f"Acteurs: {row['Actors']}"

def add_arch_overview_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT_LAYOUT])
    slide.shapes.title.text = "Architecture – Vue d'ensemble"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Consultez le diagramme PlantUML généré dans la documentation."

def generate_ppt(xlsx: pathlib.Path, output: pathlib.Path):
    ctx, project = read_ctx(xlsx)
    prs = Presentation()
    add_title_slide(prs, project)
    add_objectives_slide(prs, ctx)
    add_process_slide(prs, ctx)
    add_arch_overview_slide(prs)
    prs.save(output)
    print(f"✅ Présentation générée: {output}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument('-i', '--input', default='petstore_archi_v3.xlsx')
    parser.add_argument('-o', '--output', default='presentation.pptx')
    args = parser.parse_args()
    generate_ppt(pathlib.Path(args.input), pathlib.Path(args.output))
